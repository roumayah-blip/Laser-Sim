"""Generate a single-slide PowerPoint summary of a completed amplifier run."""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

import numpy as np

from laser_sim.constants import DB_PER_NP


@dataclass
class SlideOptions:
    title: str = "Yb Fiber Amplifier Simulation"
    subtitle: str = ""
    template_path: str | None = None
    output_path: str = "amplifier_summary.pptx"


def build_amplifier_slide(inp, outcome, options: SlideOptions) -> str:
    """
    Build a 16:9 slide with parameters, four plots, and key metrics.

    Returns path to saved ``.pptx``.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    result = outcome.result
    if result is None:
        raise ValueError("Cannot build slide: outcome.result is None (run failed).")

    plot_bufs = {
        "pump_gain_vs_z": _make_pump_gain_z_plot(result, inp),
        "signal_temporal": _make_signal_temporal_plot(result, inp),
        "output_packet": _make_output_packet_plot(result, inp),
        "b_integral": _make_b_integral_plot(outcome),
    }

    if options.template_path and Path(options.template_path).is_file():
        prs = Presentation(options.template_path)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.6))
    tf = title_box.text_frame
    tf.text = options.title or "Yb Fiber Amplifier Simulation"
    tf.paragraphs[0].runs[0].font.size = Pt(28)
    tf.paragraphs[0].runs[0].font.bold = True

    if options.subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.4))
        sub_box.text_frame.text = options.subtitle
        sub_box.text_frame.paragraphs[0].runs[0].font.size = Pt(12)

    _add_parameter_table(
        slide,
        inp,
        result,
        outcome,
        left=Inches(0.4),
        top=Inches(1.3),
        width=Inches(5.0),
        height=Inches(5.6),
    )

    plot_left = Inches(5.6)
    plot_top = Inches(1.3)
    plot_w = Inches(3.85)
    plot_h = Inches(2.7)
    slide.shapes.add_picture(plot_bufs["pump_gain_vs_z"], plot_left, plot_top, plot_w, plot_h)
    slide.shapes.add_picture(
        plot_bufs["signal_temporal"],
        plot_left + plot_w + Inches(0.05),
        plot_top,
        plot_w,
        plot_h,
    )
    slide.shapes.add_picture(
        plot_bufs["output_packet"],
        plot_left,
        plot_top + plot_h + Inches(0.1),
        plot_w,
        plot_h,
    )
    slide.shapes.add_picture(
        plot_bufs["b_integral"],
        plot_left + plot_w + Inches(0.05),
        plot_top + plot_h + Inches(0.1),
        plot_w,
        plot_h,
    )

    foot = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4))
    e_out_uj = result.energy_packet_out_j * 1e6
    e_in_uj = result.energy_packet_in_j * 1e6
    gain_db = 10.0 * np.log10(max(e_out_uj / max(e_in_uj, 1e-30), 1e-30))
    b_total = outcome.b_integral.b_total_rad if outcome.b_integral else 0.0
    sev = outcome.b_integral.severity if outcome.b_integral else "n/a"
    foot.text_frame.text = (
        f"E_in={e_in_uj:.3f} µJ  →  E_out={e_out_uj:.3f} µJ  "
        f"(Gain {gain_db:+.2f} dB, ×{e_out_uj / max(e_in_uj, 1e-30):.1f})    "
        f"B_total={b_total:.3f} rad ({sev})"
    )
    foot.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    out_path = str(options.output_path)
    prs.save(out_path)
    return out_path


def _make_pump_gain_z_plot(result, inp) -> BytesIO:
    import matplotlib.pyplot as plt

    fig, ax1 = plt.subplots(figsize=(5, 3.2), dpi=150)
    z_mm = result.z_m * 1e3
    p_pump_mean = np.mean(result.pump_fwd_w, axis=1)
    if p_pump_mean.size:
        p0 = max(float(p_pump_mean[0]), 1e-30)
        ax1.plot(z_mm, (1.0 - p_pump_mean / p0) * 100.0, color="tab:blue", lw=1.5)
        ax1.set_ylabel("Pump absorbed (%)", color="tab:blue")
        ax1.tick_params(axis="y", labelcolor="tab:blue")
    ax2 = ax1.twinx()
    t = result.t_s
    burst_start = getattr(inp, "burst_start_s", None)
    if burst_start is None and hasattr(inp, "signal_channels_inputs"):
        burst_start = 0.0
    else:
        burst_start = float(getattr(inp, "burst_start_s", 200e-6))
    it_ss = int(np.searchsorted(t, burst_start))
    it_ss = min(max(it_ss, 0), t.size - 1)
    n2_z = result.populations.n2[:, it_ss]
    n0_z = result.populations.n0[:, it_ss]
    wl_center = float(getattr(inp, "signal_center_nm", 1030.0))
    j_ctr = int(np.argmin(np.abs(result.wavelength_nm - wl_center)))
    se = result.sigma_em_signal_m2[j_ctr]
    sa = result.sigma_abs_signal_m2[j_ctr]
    gamma_s = float(getattr(result, "gamma_signal", 0.85))
    g_z = gamma_s * (se * n2_z - sa * n0_z)
    ax2.plot(z_mm, g_z * DB_PER_NP, color="tab:red", lw=1.5)
    ax2.set_ylabel("g₀ (dB/m)", color="tab:red")
    ax2.axhline(0, color="gray", lw=0.5, ls="--")
    ax1.set_xlabel("Position z (mm)")
    ax1.set_title("Pump absorption & small-signal gain")
    fig.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _make_signal_temporal_plot(result, inp) -> BytesIO:
    import matplotlib.pyplot as plt

    fig, ax1 = plt.subplots(figsize=(5, 3.2), dpi=150)
    t = result.t_s
    dlam = np.gradient(result.wavelength_nm)
    p_in_t = np.sum(result.signal_fwd_w_nm[0] * dlam, axis=1)
    p_out_t = np.sum(result.signal_fwd_w_nm[-1] * dlam, axis=1)
    chirp = float(getattr(inp, "chirp_duration_s", 1e-9))
    burst_start = float(getattr(inp, "burst_start_s", 200e-6))
    burst_n = int(getattr(inp, "burst_count", 5))
    burst_sp = float(getattr(inp, "burst_spacing_s", 2.5e-9))
    t_center = burst_start + 0.5 * chirp
    half = max(2e-9, 0.5 * chirp + max(burst_n - 1, 0) * burst_sp)
    mask = (t >= t_center - half) & (t <= t_center + half)
    t_ns = (t[mask] - t_center) * 1e9
    ax1.plot(t_ns, p_in_t[mask], color="tab:gray", lw=1.2)
    ax2 = ax1.twinx()
    ax2.plot(t_ns, p_out_t[mask], color="tab:red", lw=1.5)
    ax1.set_xlabel("Time (ns, packet-centered)")
    ax1.set_ylabel("P_in (W)", color="tab:gray")
    ax2.set_ylabel("P_out (W)", color="tab:red")
    ax1.set_title("Signal temporal (in vs out)")
    fig.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _make_output_packet_plot(result, inp) -> BytesIO:
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(figsize=(5, 3.2), dpi=150)
    pe = getattr(result, "pulse_energies_out_j", None)
    if pe is not None and len(pe) > 1:
        e_out_uj = np.asarray(pe) * 1e6
        ax.bar(range(1, len(e_out_uj) + 1), e_out_uj, color="tab:blue")
        ax.set_xlabel("Pulse index")
        ax.set_ylabel("Output energy (µJ)")
        ax.set_title(f"Packet output — {len(e_out_uj)} pulses, total {e_out_uj.sum():.2f} µJ")
    else:
        t = result.t_s
        dlam = np.gradient(result.wavelength_nm)
        p_out_t = np.sum(result.signal_fwd_w_nm[-1] * dlam, axis=1)
        chirp = float(getattr(inp, "chirp_duration_s", 1e-9))
        burst_start = float(getattr(inp, "burst_start_s", 200e-6))
        t_center = burst_start + 0.5 * chirp
        half = 3 * chirp
        mask = (t >= t_center - half) & (t <= t_center + half)
        ax.plot((t[mask] - t_center) * 1e9, p_out_t[mask], color="tab:red", lw=1.5)
        ax.set_xlabel("Time (ns)")
        ax.set_ylabel("Output power (W)")
        ax.set_title(f"Output — E={result.energy_packet_out_j * 1e6:.3f} µJ")
    fig.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _make_b_integral_plot(outcome) -> BytesIO:
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(figsize=(5, 3.2), dpi=150)
    b = outcome.b_integral
    if b is None:
        ax.text(0.5, 0.5, "B-integral not computed", ha="center", va="center", transform=ax.transAxes)
    else:
        labels = ["Pre", "Active", "Post"]
        vals = [b.b_passive_before_rad, b.b_active_rad, b.b_passive_after_rad]
        color = {
            "excellent": "tab:green",
            "moderate": "tab:olive",
            "significant": "tab:orange",
            "severe": "tab:red",
        }.get(b.severity, "tab:blue")
        ax.bar(labels, vals, color=color)
        ax.axhline(1.0, color="gray", lw=0.5, ls="--")
        ax.axhline(3.0, color="orange", lw=0.5, ls="--")
        ax.axhline(5.0, color="red", lw=0.5, ls="--")
        ax.set_ylabel("B-integral (rad)")
        ax.set_title(f"B-integral: total {b.b_total_rad:.3f} rad ({b.severity})")
        ax.legend(loc="upper left", fontsize=7)
    fig.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _add_parameter_table(slide, inp, result, outcome, *, left, top, width, height) -> None:
    from pptx.util import Pt

    rows_data = _collect_parameter_rows(inp, result, outcome)
    table = slide.shapes.add_table(len(rows_data) + 1, 2, left, top, width, height).table
    table.cell(0, 0).text = "Parameter"
    table.cell(0, 1).text = "Value"
    for cell in (table.cell(0, 0), table.cell(0, 1)):
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(10)
    for i, (k, v) in enumerate(rows_data, start=1):
        table.cell(i, 0).text = k
        table.cell(i, 1).text = v
        for col in (0, 1):
            for p in table.cell(i, col).text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)


def _collect_parameter_rows(inp, result, outcome) -> list[tuple[str, str]]:
    rows: list[tuple[str, str]] = []
    rows.append(("Material", f"{inp.material_key}"))
    rows.append(("Core diameter", f"{inp.core_diameter_um:.2f} µm"))
    rows.append(("Cladding diameter", f"{inp.cladding_diameter_um:.2f} µm"))
    rows.append(("Fiber length", f"{inp.fiber_length_m * 1000:.1f} mm"))
    rows.append(("Pumping", "Cladding" if inp.cladding_pumped else "Core"))
    if inp.abs_mode_db_per_m:
        rows.append(("Pump absorption (dB/m)", f"{inp.pump_absorption_db_per_m:.3f}"))
    rows.append(("Pump λ", f"{inp.pump_wavelength_nm:.1f} nm"))
    rows.append(("Pump power", f"{inp.pump_peak_power_w:.3f} W"))
    rows.append(("Pump mode", "CW" if inp.pump_cw else inp.pump_shape))
    rows.append(("Signal center λ", f"{inp.signal_center_nm:.2f} nm"))
    rows.append(("Packet energy", f"{inp.packet_energy_j * 1e6:.3f} µJ"))
    rows.append(("Burst count", f"{inp.burst_count}"))
    rows.append(("n_z", f"{inp.n_z}"))
    rows.append(("n_t", f"{result.t_s.size}"))
    rows.append(("n_λ", f"{result.wavelength_nm.size}"))
    rows.append(("─── Outputs ───", ""))
    rows.append(("E_in (packet)", f"{result.energy_packet_in_j * 1e6:.3f} µJ"))
    rows.append(("E_out (packet)", f"{result.energy_packet_out_j * 1e6:.3f} µJ"))
    gain_lin = result.energy_packet_out_j / max(result.energy_packet_in_j, 1e-30)
    rows.append(("Energy gain", f"×{gain_lin:.2f} ({10 * np.log10(max(gain_lin, 1e-30)):+.2f} dB)"))
    if outcome.b_integral:
        rows.append(
            (
                "B-integral",
                f"{outcome.b_integral.b_total_rad:.3f} rad ({outcome.b_integral.severity})",
            )
        )
    if getattr(inp, "multichannel_mode", False):
        rows.append(("Multi-channel", "yes"))
    return rows
